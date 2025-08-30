from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.shapes import PP_PLACEHOLDER
import io
import random

# Map our layout hints to likely layout names in real templates.
LAYOUT_PREFS = [
    ("title_and_content", ["Title and Content", "Content with Caption", "Title and Content (2)", "Title and Content 2"]),
    ("title_only",        ["Title Only", "Blank Title", "Title"]),
    ("section_header",    ["Section Header", "Section Title", "Title Slide"]),
    ("two_content",       ["Two Content", "Two Content and Title", "Comparison"]),
    ("quote",             ["Quote", "Title Only"]),
    ("comparison",        ["Comparison", "Two Content"]),
    ("timeline",          ["Title and Content", "Two Content"]),
    ("process",           ["Title and Content", "Two Content"]),
    ("overview",          ["Title and Content", "Title Only"]),
    ("summary",           ["Title and Content", "Title Only"]),
]

TITLE_TYPES = {PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE, PP_PLACEHOLDER.SUBTITLE}
CONTENT_TYPES = {PP_PLACEHOLDER.BODY, PP_PLACEHOLDER.OBJECT}  # BODY in most templates; OBJECT often has text_frame too
PICTURE_TYPES = {PP_PLACEHOLDER.PICTURE, PP_PLACEHOLDER.SLIDE_IMAGE}

def _layout_has_body(layout):
    """Return True if this layout has a BODY (or text-capable non-title) placeholder."""
    try:
        for ph in layout.placeholders:
            t = ph.placeholder_format.type
            if t in CONTENT_TYPES:
                return True
            if t not in TITLE_TYPES and hasattr(ph, "text_frame") and ph.text_frame is not None:
                return True
    except Exception:
        pass
    return False

def _find_layout_index(prs: Presentation, layout_hint: str) -> int:
    """Pick a layout index; if the hint fails, prefer any layout that actually has BODY content."""
    prefs = next((p[1] for p in LAYOUT_PREFS if p[0] == layout_hint), None) or ["Title and Content", "Title Only"]
    # Try preferred names first
    for name in prefs:
        for i, layout in enumerate(prs.slide_layouts):
            try:
                lname = getattr(layout, "name", "") or ""
                if name.lower() in lname.lower():
                    return i
            except Exception:
                pass
    # Prefer any layout that clearly has BODY/text content
    for i, layout in enumerate(prs.slide_layouts):
        if _layout_has_body(layout):
            return i
    # Fallback: any layout containing "title"
    for i, layout in enumerate(prs.slide_layouts):
        try:
            lname = getattr(layout, "name", "") or ""
            if "title" in lname.lower():
                return i
        except Exception:
            pass
    return 0

def _collect_template_images(prs: Presentation):
    """Collect (blob, width, height) of images present in the uploaded template/presentation."""
    images = []
    for slide in prs.slides:
        for shape in slide.shapes:
            try:
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE and hasattr(shape, "image"):
                    images.append((shape.image.blob, shape.width, shape.height))
            except Exception:
                continue
    random.shuffle(images)
    return images

def _purge_all_existing_slides(prs: Presentation):
    """Remove all existing slides while preserving masters/theme."""
    sldIdLst = prs.slides._sldIdLst
    for sldId in list(sldIdLst):
        rId = sldId.rId
        prs.part.drop_rel(rId)
        sldIdLst.remove(sldId)

def _first_placeholder(slide, allowed_types):
    """Return the first placeholder whose placeholder_format.type is in allowed_types."""
    for shp in slide.placeholders:
        try:
            if shp.placeholder_format and shp.placeholder_format.type in allowed_types:
                return shp
        except Exception:
            pass
    return None

def _first_text_capable_non_title(slide):
    """Find any non-title placeholder that has a text_frame (for quirky templates)."""
    for shp in slide.placeholders:
        try:
            t = shp.placeholder_format.type
            if t in TITLE_TYPES:
                continue
            if hasattr(shp, "text_frame") and shp.text_frame is not None:
                return shp
        except Exception:
            pass
    return None

def _set_text(shape, text: str):
    """Set text into a placeholder/textbox while letting the template's default run styles apply."""
    try:
        if not hasattr(shape, "text_frame") or shape.text_frame is None:
            return
        tf = shape.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = (text or "")
        # Avoid stomping on theme fonts; only shrink if extremely long
        if len(text or "") > 90:
            for r in p.runs:
                r.font.size = Pt(20)
    except Exception:
        pass

def _set_bullets(shape, bullets):
    """Write bullets into the content placeholder and keep theme paragraph styles."""
    try:
        if not hasattr(shape, "text_frame") or shape.text_frame is None:
            return
        tf = shape.text_frame
        tf.clear()
        bullets = (bullets or [])[:12]
        for i, b in enumerate(bullets):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            run = p.add_run()
            run.text = b
            p.level = 0
            if len(b) > 100:  # gentle shrink for super-long bullets
                for r in p.runs:
                    r.font.size = Pt(16)
        tf.word_wrap = True
    except Exception:
        pass

def _ensure_notes(slide, text: str):
    """Add speaker notes if requested; ignore failures if template lacks notes master."""
    if not text:
        return
    try:
        notes = slide.notes_slide
        if notes and notes.notes_text_frame:
            notes.notes_text_frame.text = text
    except Exception:
        pass

def _fill_picture_placeholder_if_any(slide, image_blob: bytes):
    """If the layout provides a picture placeholder, insert image there to keep styling."""
    try:
        pic_ph = _first_placeholder(slide, PICTURE_TYPES)
        if pic_ph is not None:
            try:
                # Some picture placeholders support .insert_picture()
                pic_ph.insert_picture(io.BytesIO(image_blob))
                return True
            except Exception:
                # fallback: draw over the placeholder rect
                left, top, width, height = pic_ph.left, pic_ph.top, pic_ph.width, pic_ph.height
                slide.shapes.add_picture(io.BytesIO(image_blob), left, top, width=width, height=height)
                return True
    except Exception:
        pass
    return False

def build_presentation(template_bytes: bytes, slides_plan):
    """
    Build a new PPTX from the uploaded template and an LLM-produced slide plan.
    slides_plan is a list of dicts: {title, bullets, layout_hint, notes?}
    """
    prs = Presentation(io.BytesIO(template_bytes))

    # 1) Collect images then purge slides so template content doesn't leak in.
    template_images = _collect_template_images(prs)
    _purge_all_existing_slides(prs)

    # 2) Build slides from plan
    for idx, slide_data in enumerate(slides_plan):
        hint = slide_data.get("layout_hint", "title_and_content")
        layout_idx = _find_layout_index(prs, hint)
        slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])

        # If the chosen layout still lacks a body placeholder, try to switch to one that has it.
        if _first_placeholder(slide, CONTENT_TYPES) is None and _first_text_capable_non_title(slide) is None:
            for i, layout in enumerate(prs.slide_layouts):
                if _layout_has_body(layout):
                    slide = prs.slides.add_slide(layout)  # add a new one with content
                    break

        title_text = (slide_data.get("title") or "")[:200]
        bullets = (slide_data.get("bullets") or [])[:12]

        # Title: prefer slide.shapes.title first (respects theme), then title placeholders.
        title_ph = getattr(slide.shapes, "title", None)
        if title_ph is None:
            title_ph = _first_placeholder(slide, TITLE_TYPES)

        # Body/content: BODY, OBJECT (text), or any non-title text-capable placeholder.
        body_ph = _first_placeholder(slide, CONTENT_TYPES) or _first_text_capable_non_title(slide)

        if title_ph:
            _set_text(title_ph, title_text)
        else:
            # last resort textbox (not ideal for theme)—only if nothing else to write title into
            try:
                tb = slide.shapes.add_textbox(Inches(1), Inches(0.7), Inches(8), Inches(1))
                _set_text(tb, title_text)
            except Exception:
                pass

        if body_ph:
            _set_bullets(body_ph, bullets)
        else:
            # last resort textbox for bullets
            try:
                tb = slide.shapes.add_textbox(Inches(1), Inches(1.7), Inches(8), Inches(4.5))
                _set_bullets(tb, bullets)
            except Exception:
                pass

        _ensure_notes(slide, slide_data.get("notes", ""))

        # Reuse a template image in a way that keeps layout look-and-feel
        if template_images and ((idx % 3 == 0) or len(bullets) <= 2):
            blob, w, h = template_images[idx % len(template_images)]
            placed = _fill_picture_placeholder_if_any(slide, blob)
            if not placed:
                # tasteful accent if no picture placeholder exists on layout
                try:
                    slide.shapes.add_picture(io.BytesIO(blob), Inches(0.4), Inches(5.1), height=Inches(1.2))
                except Exception:
                    pass

    out = io.BytesIO()
    prs.save(out)
    return out.getvalue()
